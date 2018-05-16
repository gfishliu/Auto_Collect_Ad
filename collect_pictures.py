# -*- coding: utf-8  -*-
import time
import cv2
import sys
import os
import pptx
import warnings
import requests
import re
import numpy as np
import logging
from selenium import webdriver
from pptx.util import Inches
import xlrd
import datetime
from xlutils.copy import copy
from PIL import ImageGrab
from PIL import Image
from bs4 import BeautifulSoup
from xlrd import xldate_as_tuple

# 设置全局参量以便识别
SHEET_A = 0
SHEET_B = 1
SHEET_C = 2
SHEET_D = 3
SHEET_E = 4
SHEET_F = 5
SHEET_G = 6
SHEET_H = 7
SHEET_I = 8
SHEET_J = 9
FAIL = "0.0"
SUCCESS = "Y"
NORMAL = "1.0"
imgPath = "result/"  # 图片保存位置
libPath = "lib/"  # 所用到的库及其他文件位置
jqueryName = "jquery-3.0.0.min.js"  # jquery库名称
confFile = "conf.xls"  # 配置名称
rules = {}  # 配置里的规则集合
cookies = {}
deadLine = '2018-06-01'  # 截止日期写死
logger = logging.getLogger("log")
browserWidth = 1500
browserHigth = 1000


# 检测汽车之家链接中的跳转链接 url: 汽车之家里的根据规则获取的链接   key: 检测代码  返回: 是否找到检测代码 True False
def check_url(url, key):
	if "ah_mark" not in url:  # 只遍历包含ah_mark的链接
		return False
	try:
		rsp = requests.get(url.strip('/'), allow_redirects=True)  # 请求链接
	except:
		rsp = None
	if rsp == None:  # 无响应则直接返回
		return False
	logger.info(url.strip('/').decode("utf-8"))
	logger.info("Redirects:")
	for r in rsp.history:  # 遍历响应的重定向链接中，是否包含监控代码key
		logger.info(r.url + "\n")
		if key in r.url:  # 包含则返回True
			return True
	logger.info("-------------")
	return False


# 处理汽车之家，遍历指定属性，获取重定向包含key的属性    driver: 浏览器驱动   rule: 过滤规则  key: 检测代码  返回: 广告位置坐标
def special_deal(driver, rule, key):
	input = "var links = []; outs=$(\"[" + rule + "]\").each(function(){links.push($(this).attr(\"" + rule + "\"))});return links;"  # 遍历指定规则，讲属性对应的值保存为links
	logger.info(input)
	links = driver.execute_script(input)
	logger.info("Links: " + str(len(links)))
	if links != None:
		for link in links:  # 遍历links
			# 遍历所有指定规则的属性元素 进行urlcheck 从重定向里找到检测代码key 找到的话，对相应的link做offset截图
			if check_url(link, key) == True:
				return normal_deal(driver, rule, link)  # 如果从links中找到了监测代码，那么进入normal_deal，获取link的位置
	return []


# 根据指定规则，获取元素在页面中的位置   driver: 浏览器驱动   rule: 过滤规则  key: 检测代码  返回: 广告位置坐标
def normal_deal(driver, rule, key):
	input = "outs=$(\"[" + rule + "*=\'" + key.replace("\n", "") + "\']\").offset();return outs;"  # 根据规则，获取指定元素的位置
	offset = driver.execute_script(input)
	logger.info(input)
	input = "return $(\"[" + rule + "*=\'" + key.replace("\n", "") + "\']\").parents(\"div\").width();"  # 获取元素父级div宽度
	pWidth = driver.execute_script(input)
	logger.info(input)
	input = "return $(\"[" + rule + "*=\'" + key.replace("\n", "") + "\']\").parents(\"div\").height();"  # 获取元素父级div高度
	pHight = driver.execute_script(input)
	logger.info(input)
	if offset == None or pWidth == None or pHight == None:
		return None
	adElem = {'top': offset['top'], 'left': offset['left'], 'width': pWidth, 'height': pHight}  # 返回广告位置元组信息
	return adElem


# 网页截图逻辑 url_name: 查找的url名称 goal_path: 保存路径  url: url链接    key:检测代码   save_name: 保存图片名称 driver: 浏览器驱动  rule: 过滤规则  返回:截取的图片集合
def save_Spicture_from_url(cookie_area, url_name, goal_path, url, key, save_name, driver, rule):
	logger.info("get picture from " + url + " ...")
	logger.info("target url : " + url + " ...")
	logger.info("target key : " + key + " ...")
	logger.info("target rule : " + rule + " ...")
	driver.get(url)  # 获取指定链接
	driver.add_cookie({'name': 'adip', 'value': cookie_area})
	with open(libPath + jqueryName, 'r') as jquery_js:  # 加载本地的jquery处理库
		jquery = jquery_js.read()
		driver.execute_script(jquery)
	imgPaths = []  # 图片名称集合

	logger.info(url_name)
	if u"汽车之家" in url_name:  # 对汽车之家网站做特殊处理，遍历ah_mark有关的链接，获取重定向中包含检测代码的链接，获取该链接偏移位置
		logger.info("Special_deal ...")
		adElem = special_deal(driver, rule, key)
	else:  # 获取包含检测代码的链接，获取该链接偏移位置
		logger.info("Normal_deal ...")
		adElem = normal_deal(driver, rule, key)

	# 根据位置信息，执行屏幕滚动，并进行截图
	if adElem != None and adElem != []:
		padding = browserHigth / 2
		moveStep = 0 if adElem['top'] - padding < 0 else adElem['top'] - padding  # 确定广告位置
		# 滚动屏幕到指定位置
		driver.execute_script("""
		(function () {
			var y = 0;
			var step = 50;
			window.scroll(0, 0);

			function f() {
			if (y < """ + str(moveStep) + """) {
				y += step;
				window.scroll(0, y);
				setTimeout(f, 100);
			} else {
					document.title += "scroll-done";
				}
			}

			setTimeout(f, 1000);
		})();
		""")
		for i in xrange(1000):
			if "scroll-done" in driver.title:
				break
			time.sleep(1)

		# 截取全屏
		if os.path.exists(goal_path) == False:
			os.mkdir(goal_path)

		#######全屏截图#####
		im = ImageGrab.grab()
		screenSourceImg = goal_path + '/' + save_name + ".screen.png"  # 图片路径和名称
		im.save(screenSourceImg)  # 全屏截图
		imgPaths.append(screenSourceImg)  # 将截取的图片路径保存到列表里
		#######全屏截图#####

		######挖取浏览器截图####
		browserTotalImg = goal_path + '/' + save_name + ".bak.png"
		driver.save_screenshot(browserTotalImg)  # 浏览器截图
		im = Image.open(goal_path + '/' + save_name + ".bak.png")
		iXpoi = 0 if adElem['left'] < 0 else adElem['left']
		iYpoi = adElem['top'] if moveStep == 0 else (browserHigth - adElem['height']) / 2
		im = im.crop(
			(iXpoi, 0 if iYpoi - 100 < 0 else iYpoi - 100, iXpoi + adElem['width'], iYpoi + adElem['height'] + 100))
		browserPartImg = goal_path + '/' + save_name + ".bak.crop.png"
		im.save(browserPartImg)
		######挖取浏览器截图####

		######模板匹配出红框效果###
		screenResultImg = goal_path + '/' + save_name + ".png"
		match_picture(screenSourceImg, browserPartImg, screenResultImg)
		######模板匹配出红框效果###
		os.remove(screenSourceImg)
		os.remove(browserPartImg)
		os.remove(browserTotalImg)

		print "Get Target Picture " + screenResultImg + " Success! ... \n"
	if len(imgPaths) != 0:  # 返回列表
		return imgPaths


# 红框效果 imgSourceName 源图, imgTemplateName 要匹配的红框图, saveImgName 保存图片名称
def match_picture(imgSourceName, imgTemplateName, saveImgName):
	imgS_rbg = cv2.imread(imgSourceName)
	imgS_gray = cv2.cvtColor(imgS_rbg, cv2.COLOR_BGR2GRAY)
	imgT = cv2.imread(imgTemplateName, 0)
	w, h = imgT.shape[::-1]
	res = cv2.matchTemplate(imgS_gray, imgT, cv2.TM_CCOEFF_NORMED)
	threshold = 0.88  # 匹配识别度调节
	loc = np.where(res >= threshold)
	for pt in zip(*loc[::-1]):
		cv2.rectangle(imgS_rbg, pt, (pt[0] + w, pt[1] + h), (0, 0, 255), 1)
	cv2.imwrite(saveImgName, imgS_rbg)


# 保存到ppt curAdnum：当前广告序号  totalAdnum：总广告序号   url: 网站链接 goal_path： 保存路径
def save_picture_to_ppt(saveImgNum, url,saveUrl, goal_path):
	if os.path.exists(goal_path) == False:  # 检查文件夹是否存在
		os.mkdir(goal_path)
	pptName = datetime.datetime.now().strftime("%Y%m%d") + ".pptx"
	goal_ppt = goal_path + '/' + pptName
	pptFile = pptx.Presentation(libPath + 'template.pptx')  # 从已有模板读入初始化

	count = 0
	for i in saveImgNum:
		fn = goal_path + '/' + str(i) + '.png'
		if os.path.exists(fn) == False:
			print fn + " is not exits! Please Check the picture ... "
			continue
		print fn
		slide = pptFile.slides.add_slide(pptFile.slide_layouts[5])  # 选取ppt样式
		slide.shapes.placeholders[0].text = url+": "+saveUrl[count]  # 设置文字
		img = cv2.imread(fn)
		sp = img.shape
		imgWidth = 9
		imgHigth = round(float(sp[0]) / float(sp[1]), 2) * 9
		slide.shapes.add_picture(fn, Inches(0.5), Inches(1.5), Inches(imgWidth), Inches(imgHigth))  # 设置图片
		pptFile.save(goal_ppt)
		count+=1
	if count != 0:
		print "Save Pictures to " + goal_ppt + " Success! ... \n"
	else:
		print "No Pictures saved to " + goal_ppt + " Failed! ... \n"

def check_conf(sheet1_urls, sheet2_rule, sheet3_cook):
	sheet1_urls_nrows = sheet1_urls.nrows  # 获取SHEET1行数
	sheet1_urls_ncols = sheet1_urls.ncols  # 获取SHEET1列数
	for i in range(1, sheet1_urls_nrows):  # 行遍历SHEET1
		for j in range(0, sheet1_urls_ncols):  # 列遍历SHEET1
			if j > SHEET_I:  # 如果不是处理状态列
				state = str(sheet1_urls.row_values(i)[j])
				if state != "" and state != SUCCESS and state != FAIL and state != NORMAL:
					print "Config Error! SHEET1 %d Rows %d Columns Have Error Content ..." % (i + 1, j + 1)
					#print state, NORMAL
					return False
				continue
			if j == SHEET_E:  # 如果是数or j == 字列
				if isinstance(sheet1_urls.row_values(i)[j],
							  str):  # or str(int(sheet1_urls.row_values(i)[j])).isdigit() == False:
					print "Config Error! SHEET1 %d Rows %d Columns is  Not Digit ..." % (i + 1, j + 1)
					return False
			else:
				if j == SHEET_D:   #描述无需判断
					continue
				if re.search(r'\s', sheet1_urls.row_values(i)[j]) != None:
					print "Config Error! SHEET1 %d Rows %d Columns Have Blank Line ..." % (i + 1, j + 1)
					return False

	sheet2_rule_nrows = sheet2_rule.nrows  # 获取SHEET2行总数
	for i in range(1, sheet2_rule_nrows):  # 行遍历SHEET2
		for j in range(1, SHEET_D):  # 列遍历SHEET2
			if sheet2_rule.row_values(i)[j] == "":
				print "Config Error! SHEET2 %d Rows %d Columns Have No Content ..." % (i + 1, j + 1)
				return False
			if re.search(r'\s', sheet2_rule.row_values(i)[j]) != None:
				print "Config Error! SHEET2 %d Rows %d Columns Have Blank Line ..." % (i + 1, j + 1)
				return False

	sheet3_cook_nrows = sheet3_cook.nrows  # 获取SHEET3行总数
	for i in range(1, sheet3_cook_nrows):  # 行遍历SHEET3
		for j in range(0, SHEET_B):  # 列遍历SHEET3
			if sheet3_cook.row_values(i)[j] == "":
				print "Config Error! SHEET3 %d Rows %d Columns Have No Content ..." % (i + 1, j + 1)
				return False
			if re.search(r'\s', sheet3_cook.row_values(i)[j]) != None:
				print "Config Error! SHEET3 %d Rows %d Columns Have Blank Line ..." % (i + 1, j + 1)
				return False

	return True


# 初始化日志模块配置
def init_log():
	logger.setLevel(level=logging.INFO)
	handler = logging.FileHandler("Running.log.txt")
	handler.setLevel(logging.INFO)
	formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
	handler.setFormatter(formatter)
	logger.addHandler(handler)


# 初始化目录
def init_dir():
	if os.path.exists(imgPath) == False:
		os.mkdir(imgPath)


def check_time():
	try:
		rsp = requests.get("http://time.tianqi.com/", allow_redirects=True)  # 请求链接
	except:
		print "Check Time! Please Check The Network ... \n"
		exit_program()

	if rsp.status_code == 200:
		soup = BeautifulSoup(rsp.text, "lxml")
		timeStr = soup.find_all(name='meta')[2].encode('gb2312')
		timeFormat = timeStr[timeStr.find("2"):timeStr.find("2") + 10]
		return timeFormat <= deadLine
	return False


def get_date_rcol(sheet1_urls):
	curDate = datetime.datetime.now().strftime("%m/%d")
	date_rcol = 0
	sheet1_urls_ncols = sheet1_urls.ncols  # 获取SHEET1列数
	for i in range(SHEET_J, sheet1_urls_ncols):  # 列遍历SHEET1
		if curDate == sheet1_urls.cell_value(0, i):
			date_rcol = i
			break
	return date_rcol

def check_sheet2_cheet3_value(sheet1_urls, rules, cookies):
	sheet1_urls_nrows = sheet1_urls.nrows
	for i in range(1, sheet1_urls_nrows):
		try:
			rule = rules[sheet1_urls.row_values(i)[SHEET_B] + sheet1_urls.row_values(i)[SHEET_I]]  # 从rules[] 读取过滤的规则
		except:
			print "Conf Error! ", sheet1_urls.row_values(i)[SHEET_B],"Rule ",sheet1_urls.row_values(i)[SHEET_B] + sheet1_urls.row_values(i)[SHEET_I]," Please Check Sheet2 Have Rule or Not  ..."
			exit_program()

		area = sheet1_urls.row_values(i)[SHEET_F]
		try:
			cookie = cookies[area]
		except:
			print "Conf Error! ", sheet1_urls.row_values(i)[SHEET_B]," Cookie ",area," Please Check Sheet3 Have Rule or Not  ..."
			exit_program()

# 主函数
def main():
	print "Program begin running ... \n"
	if check_time() == False:
		print "Check Time! Software Use Time Has Ended ... \n"
		exit_program()

	init_log()
	init_dir()
	print "Checking Conf ...\n"
	# 读取sheet2 中规则， 网址和类型作为主键
	# 遍历sheet1 根据网址和类型，匹配sheet2的规则
	try:
		book = xlrd.open_workbook(confFile)
	except:
		print "Config Error! No config.xls Exits ... \n"
		exit_program()
	wb = copy(book)
	ws = wb.get_sheet(0)
	sheet1_urls = book.sheet_by_index(0)  # 获取sheet1用例
	sheet2_rule = book.sheet_by_index(1)  # 获取sheet2用例
	sheet3_cook = book.sheet_by_index(2)  # 获取sheet3用例
	# 检测conf是否有问题
	if check_conf(sheet1_urls, sheet2_rule, sheet3_cook) == False:
		exit_program()

	sheet2_rule_nrows = sheet2_rule.nrows  # 获取行总数
	for i in range(1, sheet2_rule_nrows):  # 遍历excel列表存储到rules[]
		rules[sheet2_rule.row_values(i)[1] + sheet2_rule.row_values(i)[2]] = sheet2_rule.row_values(i)[3]
	sheet3_cook_nrows = sheet3_cook.nrows
	for i in range(1, sheet3_cook_nrows):
		cookies[sheet3_cook.row_values(i)[0]] = str(sheet3_cook.row_values(i)[1])

	check_sheet2_cheet3_value(sheet1_urls, rules, cookies)
	#sys.exit(0)
	sheet1_urls_nrows = sheet1_urls.nrows  # 获取sheet1行数
	date_rcol_num = get_date_rcol(sheet1_urls)
	if date_rcol_num == 0:
		print "Check Date! Please Check The Sheet1 Date of the Conf ... \n"
		exit_program()

	print "Check Conf Success! ...\n"

	try:
		driver = webdriver.Firefox()  # 调用Firfox API驱动
	except:
		print "Firfox Error! Check geckodriver.exe or Firefox update to Version 55  ..."
		exit_program()
	driver.set_window_size(browserWidth, browserHigth)  # 设置Firfox 窗口大小

	print "============================ Begin Deal ============================= \n"
	imgCounter = 1
	tmpAdname = "begin"
	saveUrl = []
	saveImgNum = []
	for i in range(1, sheet1_urls_nrows):  # 遍历每条待测链接
		state = str(sheet1_urls.row_values(i)[date_rcol_num])  # 处理状态
		url = sheet1_urls.row_values(i)[SHEET_H]  # 广告链接
		key = sheet1_urls.row_values(i)[SHEET_G]  # 监测代码
		adname = sheet1_urls.row_values(i)[SHEET_A]  # 广告商名
		urlname = sheet1_urls.row_values(i)[SHEET_B]  # 网站名字
		saveUrl.append(url)


		refreshNum = sheet1_urls.row_values(i)[SHEET_E]
		area = sheet1_urls.row_values(i)[SHEET_F]  # 广告链接
		#imgCaptureName = str(int(imgCounter))  # 图片名称 "序号id"
		if tmpAdname != "begin" and tmpAdname != adname:  # 如果不为空 并检查是否为广告商最后一条
			#print tmpAdname,adname
			goal_path = imgPath + '/' + tmpAdname.encode('gb2312')  # 要保存到的目标目录
			save_picture_to_ppt(saveImgNum, urlname,saveUrl, goal_path)  # 保存到ppt
			saveImgNum = []
			saveUrl = []
			imgCounter = 1
			print "======================== " + tmpAdname + " Deal End ======================== \n\n\n"
		#print tmpAdname,adname
		tmpAdname = adname
		goal_path = imgPath + '/' + tmpAdname.encode('gb2312')  # 要保存到的目标目录
		if state == "":
			imgCounter += 1
			if (i == sheet1_urls_nrows - 1):
				save_picture_to_ppt(saveImgNum, urlname, saveUrl, goal_path)
				print "======================== " + tmpAdname + " Deal End ======================== \n\n\n"
			continue
		elif state == SUCCESS:
			saveImgNum.append(imgCounter)
			saveUrl.append(url)
			imgCounter += 1
			if (i == sheet1_urls_nrows - 1):
				save_picture_to_ppt(saveImgNum, urlname, saveUrl, goal_path)
				print "======================== " + tmpAdname + " Deal End ======================== \n\n\n"
			continue
		elif state == FAIL:
			imgCounter += 1
			if (i == sheet1_urls_nrows - 1):
				save_picture_to_ppt(saveImgNum, urlname, saveUrl, goal_path)
				print "======================== " + tmpAdname + " Deal End ======================== \n\n\n"
			continue
		elif state == NORMAL:
			pass
		imgCaptureName = str(int(imgCounter))  # 图片名称 "序号id"
		for j in range(int(refreshNum)):
			rule = rules[sheet1_urls.row_values(i)[SHEET_B] + sheet1_urls.row_values(i)[SHEET_I]]
			#try:
			print "Refresh Total Times : %d , This is The %d Times ...\n" % (int(refreshNum), j + 1)
			imgResultPicture = save_Spicture_from_url(cookies[area], urlname, goal_path, url.replace("\n", ""), key.replace("\n", ""),
													  imgCaptureName, driver, rule)  # 处理过程
			#except:
			#	print "Firefox Error! Firfox is not in Running ..."
			#	exit_program()
			if imgResultPicture != None:
				try:
					ws.write(i, date_rcol_num, SUCCESS)  # 设置 是否完成为True
					wb.save(confFile)
					saveImgNum.append(imgCounter)
					saveUrl.append(url)
				except:
					print "Config Error! Please Close config.xls ..."
					exit_program()
				break;
			else:
				ws.write(i, date_rcol_num, FAIL)  # 设置 是否完成为False
				wb.save(confFile)
		imgCounter += 1
		if (i == sheet1_urls_nrows-1):
			save_picture_to_ppt(saveImgNum,  urlname,saveUrl, goal_path)
			print "======================== " + tmpAdname + " Deal End ======================== \n\n\n"

	driver.quit()  # 关闭驱动

def exit_program():
	print "\n"
	secs=10
	while secs:
		print "Program will exit in %d sec ..." % (secs)
		secs -= 1
		time.sleep(1)
	sys.exit(1)




if __name__ == "__main__":
	warnings.filterwarnings('ignore')  # 忽略warning
	print "\n"
	main()  # 主函数
	os.system('pause')  # 处理结束后暂停
